#!/usr/bin/env python3
"""Build the CareGrid model-overview PowerPoint from the live screenshots."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
ASSETS = HERE / "assets"
OUT = HERE / "CareGrid_Model_Overview.pptx"

EMU_IN = 914400
SW, SH = 13.333, 7.5

NAVY = RGBColor(0x12, 0x35, 0x5B)
NAVY_DK = RGBColor(0x0B, 0x22, 0x3C)
INK = RGBColor(0x17, 0x21, 0x2B)
MUTED = RGBColor(0x5D, 0x6B, 0x78)
BORDER = RGBColor(0xD9, 0xE0, 0xE6)
BG = RGBColor(0xF4, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SEV = RGBColor(0xD6, 0x45, 0x45)
SURV = RGBColor(0x2F, 0x9E, 0x62)
WAIT = RGBColor(0x2F, 0x6F, 0xD0)
GOLD = RGBColor(0xC1, 0x95, 0x2D)
LIGHT_GOLD = RGBColor(0xFD, 0xF6, 0xE3)

prs = Presentation()
prs.slide_width = int(EMU_IN * SW)
prs.slide_height = int(EMU_IN * SH)
BLANK = prs.slide_layouts[6]
FONT = "Avenir Next"


def add_slide() -> object:
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, fill, line=None, radius=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    if radius is not None:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    sp.shadow.inherit = False
    return sp


def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def para(tf, text, size=14, color=INK, bold=False, italic=False, first=False,
         bullet=None, space_after=6, align=PP_ALIGN.LEFT, font=FONT, line_spacing=1.0):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    if line_spacing:
        p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = font
    return p


def kicker_title(slide, kicker, title, title_size=30, x=0.55, y=0.42):
    kf = textbox(slide, x, y + 0.02, 11.5, 0.3)
    para(kf, kicker.upper(), size=12, color=SURV, bold=True, first=True, space_after=0)
    tf = textbox(slide, x, y + 0.30, 12.2, 0.75)
    para(tf, title, size=title_size, color=NAVY, bold=True, first=True, space_after=0)


def page_footer(slide, n):
    tf = textbox(slide, 12.35, 7.08, 0.9, 0.35)
    para(tf, str(n), size=11, color=MUTED, first=True, align=PP_ALIGN.RIGHT, space_after=0)


def shot_slide(n, num, asset, kicker, title, callouts, foot=None):
    """A screenshot slide with a gutter of UNIQUE/HOW/BENEFIT callout cards."""
    slide = add_slide()
    rect(slide, 0, 0, SW, SH, BG)
    rect(slide, 0, 0, SW, 0.14, NAVY)
    kicker_title(slide, kicker, title)
    page_footer(slide, num)

    path = ASSETS / asset
    from PIL import Image as _I
    iw, ih = _I.open(path).size
    shot_x, shot_y = 0.55, 1.42
    shot_w = 8.55
    shot_h = shot_w * ih / iw
    avail_h = 7.03 - shot_y
    if shot_h > avail_h:
        shot_h = avail_h
        shot_w = shot_h * iw / ih
    rect(slide, shot_x - 0.08, shot_y - 0.08, shot_w + 0.16, shot_h + 0.16, WHITE,
         line=BORDER, radius=0.04)
    slide.shapes.add_picture(str(path), Inches(shot_x), Inches(shot_y), width=Inches(shot_w))

    gx, gw = 9.42, 3.40
    gy = shot_y
    ch = 1.78
    gap = 0.20
    for i, (tag, tag_color, ctitle, body) in enumerate(callouts):
        cy = gy + i * (ch + gap)
        card = rect(slide, gx, cy, gw, ch, WHITE, line=BORDER, radius=0.06)
        rect(slide, gx, cy, 0.10, ch, tag_color)
        tf = textbox(slide, gx + 0.24, cy + 0.14, gw - 0.45, ch - 0.28)
        p = para(tf, tag, size=10, color=tag_color, bold=True, first=True, space_after=2, font=FONT)
        para(tf, ctitle, size=13.5, color=INK, bold=True, space_after=4)
        para(tf, body, size=11.5, color=MUTED, space_after=0)
    if foot:
        ftf = textbox(slide, 0.55, 7.06, 9.0, 0.38)
        para(ftf, foot, size=12, color=NAVY, italic=True, first=True, space_after=0)
    return slide


def bullet_card(slide, x, y, w, h, title, body, accent):
    rect(slide, x, y, w, h, WHITE, line=BORDER, radius=0.07)
    rect(slide, x, y, 0.09, h, accent)
    tf = textbox(slide, x + 0.22, y + 0.16, w - 0.42, h - 0.32)
    para(tf, title, size=13, color=NAVY, bold=True, first=True, space_after=4)
    para(tf, body, size=11.5, color=MUTED, space_after=0)


# --------------------------------------------------------------------------- 1 title
slide = add_slide()
rect(slide, 0, 0, SW, SH, NAVY)
rect(slide, 0, 0, SW, 0.16, SURV)
rect(slide, 0.7, 5.9, 11.9, 0.02, NAVY_DK)
tf = textbox(slide, 0.9, 2.6, 11.5, 1.2)
para(tf, "LIVE FROM THE RUNNING SYSTEM", size=13, color=SURV, bold=True, first=True, space_after=10)
para(tf, "CareGrid", size=54, color=WHITE, bold=True, space_after=6)
para(tf, "ICU Bed-Arbitration Decision Support", size=22, color=RGBColor(0xCF, 0xDE, 0xEE), space_after=0)
tf = textbox(slide, 0.9, 4.15, 11.5, 0.9)
para(tf, "Scoring the critical-care waitlist on severity, survival likelihood, and waiting time — "
         "explained at every layer, audited on every re-rank, and decided by the clinician.",
     size=15, color=WHITE, first=True, line_spacing=1.15, space_after=0)
metrics = [
    ("XGBoost survival model", "SOFA + age + comorbidity"),
    ("AUC 0.61 · Brier 0.121", "hold-out validated before use"),
    ("3 weight profiles", "policy stance as a named preset"),
    ("SHAP-explainable", "every prediction attributed"),
    ("Clinician in the loop", "confirm or deviate, always logged"),
]
mx = 0.9
mw, mh = 2.32, 1.06
mgg = 0.055
for label, sub in metrics:
    rect(slide, mx, 6.35, mw, mh, NAVY_DK, line=RGBColor(0x2A, 0x4D, 0x78))
    mf = textbox(slide, mx + 0.14, 6.44, mw - 0.28, 0.9)
    para(mf, label, size=12, color=WHITE, bold=True, first=True, space_after=2)
    para(mf, sub, size=9.5, color=RGBColor(0x9F, 0xB4, 0xCC), space_after=0)
    mx += mw + mgg

# --------------------------------------------------------------------------- 2 problem
n = 2
slide = add_slide()
rect(slide, 0, 0, SW, SH, BG)
rect(slide, 0, 0, SW, 0.14, NAVY)
kicker_title(slide, "The problem", "Which waiting patient gets the next ICU bed?")
tf = textbox(slide, 0.55, 1.28, 12.3, 0.9)
para(tf, "Critical-care capacity is fixed and the queue is never empty. Every candidate — the sickest, "
         "the most likely to survive, the longest-waiter — is a real human whose claim the team must weigh "
         "under time pressure and public scrutiny.",
     size=15, color=INK, first=True, line_spacing=1.2, space_after=0)
cards = [
    ("Zero-sum scarcity", "Every occupied ICU bed blocks another patient; choosing is inevitable, so the "
     "choice must be defensible.", SEV),
    ("Conflicting signals", "Severity, survival likelihood, and time waited pull in opposite directions — "
     "the ranking has to make the trade-off explicit.", WAIT),
    ("Opaque scoring", "Existing triage scores rank the queue but never say why — no factor breakdown, "
     "no explanation a nurse can challenge.", NAVY),
    ("Unreviewable decisions", "Allocations decided in a human's head leave no record — no replay, no "
     "evidence for a morbidity-and-mortality review.", SURV),
]
for i, (t, b, a) in enumerate(cards):
    col = i % 2
    row = i // 2
    bullet_card(slide, 0.55 + col * 6.27, 2.55 + row * 2.15, 6.05, 1.95, t, b, a)
tf = textbox(slide, 0.55, 6.85, 12.3, 0.4)
para(tf, "What if the queue could rank itself — and be able to say why?",
     size=15, color=NAVY, bold=True, first=True, space_after=0)
page_footer(slide, n)

# --------------------------------------------------------------------------- 3 model stack
n = 3
slide = add_slide()
rect(slide, 0, 0, SW, SH, BG)
rect(slide, 0, 0, SW, 0.14, NAVY)
kicker_title(slide, "The model", "One queue, three signals — explained at every layer")
steps = [
    ("1 · Clinical inputs", "Six SOFA organ components + age + comorbidity — quantities a clinician "
     "already knows and can verify on the chart.", SEV),
    ("2 · Survival model", "CPU gradient-boosted trees trained on 2,880 ICU-patient outcomes; predicts "
     "each candidate's survival probability per episode.", WAIT),
    ("3 · Validated gate", "Hold-out AUC 0.61 · Brier 0.121 · decile calibration — a demo refuses to run "
     "on a model that fails the tolerance.", GOLD),
    ("4 · Explained rank", "Priority Score = w₁·severity + w₂·survival + w₃·wait with SHAP attribution on "
     "the survival term, a deterministic Tie-Break Cascade, and a snapshot audit trail.", SURV),
]
sw_, sh_ = 5.98, 1.9
for i, (t, b, a) in enumerate(steps):
    col = i % 2
    row = i // 2
    x = 0.55 + col * (sw_ + 0.24)
    y = 1.6 + row * (sh_ + 0.24)
    rect(slide, x, y, sw_, sh_, WHITE, line=BORDER, radius=0.06)
    rect(slide, x, y, sw_, 0.07, a)
    tf = textbox(slide, x + 0.22, y + 0.2, sw_ - 0.44, sh_ - 0.34)
    para(tf, t, size=15, color=NAVY, bold=True, first=True, space_after=5)
    para(tf, b, size=12, color=MUTED, space_after=0)
tf = textbox(slide, 0.55, 6.0, 12.3, 0.9)
para(tf, "The scoring model combines a trained ML survival predictor with transparent, deterministic "
         "arbitration logic — so the machine-learning part is explainable and the fairness part is "
         "reproducible.",
     size=14.5, color=INK, first=True, line_spacing=1.2, space_after=0)
page_footer(slide, n)

# --------------------------------------------------------------------------- 4 live queue (01)
shot_slide(
    4, 4, "01-live-queue.png", "Live dashboard — the ranked queue",
    "Every priority score you can trace",
    [
        ("WHY IT'S UNIQUE", NAVY, "Transparent scoring",
         "Each score is the weighted sum 0.582 = 0.5·sev 0.458 + 0.3·surv 0.916 + 0.2·wait 0.391 under "
         "the 50/30/20 profile — with the three normalized factors drawn as bars beside the number."),
        ("WHY IT'S UNIQUE", NAVY, "SHAP on every prediction",
         "Green/red chips attribute each survival probability: the SOFA-20 patient's −3.80 sofa term and "
         "+1.04 age term explain its 0.284 p-surv at a glance."),
        ("KEY BENEFIT", SURV, "Defensible to the bedside",
         "Every rank maps to medical quantities — six SOFA organs, age, comorbidity, hours waited — a "
         "clinician can verify rather than trust."),
    ],
    foot="Live view: Severity-dominant (50/30/20) profile · 24h wait horizon · as of 2026-01-01 11:00.",
)

# --------------------------------------------------------------------------- 5 history (02)
shot_slide(
    5, 5, "02-history-expanded.png", "Live dashboard — per-patient traceability",
    "A journey for every patient in the queue",
    [
        ("WHY IT'S UNIQUE", NAVY, "Rank history per patient",
         "Click any row and the patient's rank on every snapshot unfolds, each chip labelled with the "
         "trigger — arrival, removal, freed bed — behind that move."),
        ("WHY IT'S UNIQUE", NAVY, "Movement signs",
         "Up, down, unchanged, new: the live queue wears its own history on its sleeve as the ward "
         "changes around it."),
        ("KEY BENEFIT", SURV, "Accountability",
         "\u201cWhy did this patient's priority change?\u201d has a documented, replayable answer — "
         "not a guess."),
    ],
)

# --------------------------------------------------------------------------- 6 snapshot replay (03)
shot_slide(
    6, 6, "03-snapshot-tip-arrival.png", "Live dashboard — replay the audit trail",
    "The ward changed; the queue reacted — and it's recorded",
    [
        ("WHY IT'S UNIQUE", NAVY, "Snapshot per re-rank",
         "Every arrival, removal, or freed bed writes an immutable Ranking Snapshot. The header shows "
         "record #3 — re-rank: tip-arrival — with the newly arrived SOFA-20 patient in the mix."),
        ("HOW IT WORKS", WAIT, "Waiting time saturates",
         "The wait factor follows a 24h quadratic horizon then caps — protecting the 36h long-waiter ("
         "patient-1) without ever letting a wait dominate the decision."),
        ("KEY BENEFIT", SURV, "Reproduce any instant",
         "Drag the slider back in time: the queue at any recorded moment is exactly what was ranked "
         "then — recreatable, reviewable, shareable."),
    ],
)

# --------------------------------------------------------------------------- 7 decision (04)
shot_slide(
    7, 7, "04-decision-card.png", "Live dashboard — the arbitration decision",
    "The system recommends. The clinician decides.",
    [
        ("HOW IT WORKS", WAIT, "Recommendation with reasoning",
         "When a bed frees, the panel states the case: \u201cpatient-1 is ranked #1 with priority score "
         "0.673 — severity 0.417, survival 0.8…\u201d — the evidence, next to the queue."),
        ("WHY IT'S UNIQUE", NAVY, "Confirm or consciously deviate",
         "The clinician confirms — or overrides to a lower-ranked candidate, and the deviation, the "
         "recommendation it overrode, and any note are recorded in the same audit trail."),
        ("KEY BENEFIT", SURV, "Human-in-the-loop by design",
         "The system never auto-allocates. The final call is always the clinician's, and always "
         "auditable after the fact."),
    ],
)

# --------------------------------------------------------------------------- 8 policy comparison (09)
shot_slide(
    8, 8, "09-cli-comparison.png", "Same queue, three policy stances — CLI replay",
    "The weight profile is a policy decision, made visible",
    [
        ("WHY IT'S UNIQUE", NAVY, "One queue, three answers",
         "The identical ward is replayed under Severity-dominant, Balanced, and Severity-heavy — and the "
         "highlighted block shows the SOFA-20 arrival overtaking to #2 when the policy says severity "
         "matters most."),
        ("HOW IT WORKS", WAIT, "Named weight presets",
         "50/30/20, 40/30/30, 60/25/15 — each Weight Profile is scoped to the run and logged with every "
         "snapshot, so a score is always traceable to its stance."),
        ("KEY BENEFIT", SURV, "Policy discussions, not black boxes",
         "Hospitals argue about weights, not about hidden logic — the ethical stance is an explicit, "
         "comparable parameter."),
    ],
)

# --------------------------------------------------------------------------- 9 model gate (07)
shot_slide(
    9, 9, "07-cli-model.png", "The survival model — validated gate",
    "A demo refuses to run on an unvalidated model",
    [
        ("WHY IT'S UNIQUE", NAVY, "Hold-out validation gate",
         "AUC 0.61, Brier 0.121, per-decile calibration to 0.05 — and if the model fails the tolerance "
         "(AUC ≥ 0.60, mean calibration ≤ 0.06), every demonstration path refuses to run."),
        ("HOW IT WORKS", WAIT, "Honest metrics, in the open",
         "The report is written alongside the trained booster, decile by decile, so the strengths — and "
         "the calibration tail — are visible before anyone trusts a rank."),
        ("KEY BENEFIT", SURV, "Safe by construction",
         "You cannot demo, serve, or replay on a model that hasn't proven itself on data it never "
         "saw — the gate is in the code path, not in a review meeting."),
    ],
)

# --------------------------------------------------------------------------- 10 audit trail (08)
shot_slide(
    10, 10, "08-cli-trail.png", "End-to-end Simulation Run — the CLI trail",
    "One replayable audit trail, from ward opening to allocation",
    [
        ("HOW IT WORKS", WAIT, "Score drift you can read",
         "patient-1 climbs 0.586 → 0.673 purely as hours-on-list grow past the horizon — the waiting "
         "factor at work, snapshot by snapshot."),
        ("WHY IT'S UNIQUE", NAVY, "Append-only Ranking Snapshots",
         "Each snapshot records the ordered queue, factor breakdown, SHAP attribution, weight profile, "
         "horizon, and trigger — one continuous, reproducible story."),
        ("KEY BENEFIT", SURV, "Ready for review",
         "A morbidity-and-mortality committee can replay tonight's decision morning-after, exactly as "
         "made — including any deviation and the note."),
    ],
)

# --------------------------------------------------------------------------- 11 key benefits
n = 11
slide = add_slide()
rect(slide, 0, 0, SW, SH, NAVY)
rect(slide, 0, 0, SW, 0.16, SURV)
kicker_title(slide, "The model, in one page", "Key benefits of the CareGrid model")
benefits = [
    ("Explainable ML", "Every survival prediction carries SHAP attribution; every score decomposes into "
     "severity, survival, and waiting so a nurse can follow the math.", SURV),
    ("Deterministic fairness", "The Tie-Break Cascade resolves near-equal scores on severity → survival → "
     "wait → order-of-arrival, and says which stage decided.", GOLD),
    ("Clinician in the loop", "Recommendation always faces a human confirm-or-deviate decision, with the "
     "deviation recorded alongside the recommendation it overrode.", WAIT),
    ("Immutable audit trail", "An append-only Ranking Snapshot after every re-rank — profile, horizon, "
     "factors, SHAP, trigger — replayable end-to-end.", NAVY),
    ("Policy as configuration", "Weight Profiles are named, scoped to the run, and logged with every "
     "snapshot, so an ethical stance is explicit and comparable.", SEV),
    ("Safe by construction", "Hold-out validation gates every demonstration: no demo, dashboard, or replay "
     "runs on a model that fails AUC and calibration tolerance.", SURV),
]
bw_, bh = 5.98, 1.82
for i, (t, b, a) in enumerate(benefits):
    col = i % 2
    row = i // 2
    x = 0.55 + col * (bw_ + 0.24)
    y = 1.6 + row * (bh + 0.22)
    rect(slide, x, y, bw_, bh, NAVY_DK, line=RGBColor(0x2A, 0x4D, 0x78), radius=0.06)
    rect(slide, x, y, 0.09, bh, a)
    tf = textbox(slide, x + 0.24, y + 0.18, bw_ - 0.48, bh - 0.36)
    para(tf, t, size=15, color=WHITE, bold=True, first=True, space_after=5)
    para(tf, b, size=11.5, color=RGBColor(0xC9, 0xD8, 0xEA), space_after=0)
page_footer(slide, n)

# --------------------------------------------------------------------------- 12 close
n = 12
slide = add_slide()
rect(slide, 0, 0, SW, SH, BG)
rect(slide, 0, 0, SW, 0.14, NAVY)
tf = textbox(slide, 0.9, 2.5, 11.5, 2.4)
para(tf, "From queue to decision, explained and on record.", size=34, color=NAVY, bold=True,
     first=True, space_after=14)
para(tf, "CareGrid ranks the critical-care waitlist with a validated, explainable survival model; "
         "re-ranks on every change; and hands the clinician a recommendation on a bed they can trace, "
         "question, or override — every step written to an immutable audit trail.",
     size=16, color=MUTED, line_spacing=1.25, space_after=0)
tf = textbox(slide, 0.9, 5.6, 11.5, 0.6)
para(tf, "Not a black box. Not an autopilot. A decision-support system built to be reviewed.",
     size=15, color=SURV, bold=True, first=True, space_after=0)
page_footer(slide, n)

prs.save(OUT)
print(f"saved {OUT} with {len(prs.slides.slides if hasattr(prs.slides,'slides') else prs.slides._sldIdLst)} slides")